#!/usr/bin/env python3
"""Second round of edits to lightning-talk.pptx.

Run with: uv run --with "python-pptx pillow" modify_lightning_talk_v2.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image
import io, os

BASE    = os.path.dirname(os.path.abspath(__file__))
PPTX    = os.path.join(BASE, "lightning-talk.pptx")
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY      = RGBColor(0x1F, 0x36, 0x63)
BLUE      = RGBColor(0x00, 0x72, 0xBD)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)


# ── helpers ──────────────────────────────────────────────────────────────────

def remove_shape(slide, shape):
    shape._element.getparent().remove(shape._element)


def clear_tf(tf):
    """Remove all paragraphs after the first; remove all runs from para 0."""
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in p0._p.findall(qn('a:r')):
        p0._p.remove(r)
    return p0


def set_bold_text(shape, text, size=Pt(20), color=NAVY, align=PP_ALIGN.LEFT,
                  v_anchor=MSO_ANCHOR.MIDDLE, margin_left=Inches(0.15)):
    """Replace shape text with a single bold run."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = v_anchor
    tf.margin_left = int(margin_left)
    p0 = clear_tf(tf)
    p0.alignment = align
    run = p0.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = color


def get_blob_size(blob):
    """Return (width_px, height_px) of an image blob."""
    with Image.open(io.BytesIO(blob)) as img:
        return img.size


def get_file_size(path):
    with Image.open(path) as img:
        return img.size


def fit_in_box(img_w, img_h, box_w, box_h):
    """Return (w, h) in EMU that fills the box while preserving aspect ratio."""
    ratio = img_w / img_h
    w = box_w
    h = int(w / ratio)
    if h > box_h:
        h = box_h
        w = int(h * ratio)
    return w, h


# ── slide 2: make embedded images much larger ────────────────────────────────

def modify_slide2_images(slide):
    MARGIN = Inches(0.25)
    GAP    = Inches(0.4)
    TOP    = Inches(1.85)
    BOT    = Inches(7.35)

    box_w = int((SLIDE_W - 2 * MARGIN - GAP) / 2)
    box_h = int(BOT - TOP)

    for shape in slide.shapes:
        if shape.name == "Picture 8":
            pw, ph = get_blob_size(shape.image.blob)
            fw, fh = fit_in_box(pw, ph, box_w, box_h)
            shape.left   = int(MARGIN + (box_w - fw) // 2)
            shape.top    = int(TOP    + (box_h - fh) // 2)
            shape.width  = fw
            shape.height = fh
            print(f"    Picture 8: {fw/914400:.2f}\" × {fh/914400:.2f}\"")

        elif shape.name == "Picture 9":
            pw, ph = get_blob_size(shape.image.blob)
            fw, fh = fit_in_box(pw, ph, box_w, box_h)
            shape.left   = int(MARGIN + box_w + GAP + (box_w - fw) // 2)
            shape.top    = int(TOP    + (box_h - fh) // 2)
            shape.width  = fw
            shape.height = fh
            print(f"    Picture 9: {fw/914400:.2f}\" × {fh/914400:.2f}\"")


# ── slide 4: two side-by-side images ─────────────────────────────────────────

def modify_slide4_images(slide):
    # Remove current image and code-box
    to_remove = [s for s in slide.shapes if s.name in ("Picture 7", "Rectangle 8")]
    for s in to_remove:
        remove_shape(slide, s)
    print(f"    Removed {len(to_remove)} shapes")

    MARGIN = Inches(0.3)
    GAP    = Inches(0.4)
    TOP    = Inches(1.25)
    BOT    = Inches(7.4)

    box_w = int((SLIDE_W - 2 * MARGIN - GAP) / 2)
    box_h = int(BOT - TOP)

    # Left: rescale_v4_code.png  (portrait 0.841)
    code_path = os.path.join(BASE, "lightningtalkimages", "rescale_v4_code.png")
    cw, ch = get_file_size(code_path)
    fw, fh = fit_in_box(cw, ch, box_w, box_h)
    lx = int(MARGIN + (box_w - fw) // 2)
    ly = int(TOP    + (box_h - fh) // 2)
    slide.shapes.add_picture(code_path, lx, ly, fw, fh)
    print(f"    Left image: {fw/914400:.2f}\" × {fh/914400:.2f}\"")

    # Right: images/rescale_v4.png  (portrait 0.623)
    render_path = os.path.join(BASE, "images", "rescale_v4.png")
    rw, rh = get_file_size(render_path)
    fw2, fh2 = fit_in_box(rw, rh, box_w, box_h)
    rx = int(MARGIN + box_w + GAP + (box_w - fw2) // 2)
    ry = int(TOP    + (box_h - fh2) // 2)
    slide.shapes.add_picture(render_path, rx, ry, fw2, fh2)
    print(f"    Right image: {fw2/914400:.2f}\" × {fh2/914400:.2f}\"")


# ── slide 6: fix arrow overlapping "shared structure" text ───────────────────

def modify_slide6_arrow(slide):
    # Stack: text first (stays at ~3.17"), then arrow below text, then base class, etc.
    for shape in slide.shapes:
        name = shape.name
        if name == "TextBox 11":       # ↓ arrow — move below the text line
            shape.top = int(Inches(3.60))
        elif name == "Rectangle 13":   # ConfigurationData base class box
            shape.top = int(Inches(4.25))
        elif name == "Rectangle 14":   # 4-format insight green box
            shape.top = int(Inches(4.95))
        elif name == "TextBox 15":     # bottom punchline
            shape.top = int(Inches(5.90))
    print("    Fixed vertical stacking of arrow + text")


# ── slide 7: layout + text/title updates ─────────────────────────────────────

def modify_slide7_layout(slide):
    # Find right question box to get its top/height
    right_box = next(s for s in slide.shapes if s.name == "Rectangle 8")
    r_top = right_box.top
    r_h   = right_box.height
    r_bot = r_top + r_h

    # 4 left boxes sorted top-to-bottom
    box_names   = ("Rectangle 3", "Rectangle 4", "Rectangle 5", "Rectangle 6")
    new_labels  = (
        "Instant benchmarking",
        "Rapid prototyping",
        "Keep track of decision-making",
        "Tackle the big challenges",
    )
    left_boxes  = sorted(
        [s for s in slide.shapes if s.name in box_names],
        key=lambda s: s.top
    )

    n   = len(left_boxes)
    GAP = int(Inches(0.1))
    bh  = int((r_h - (n - 1) * GAP) / n)

    for i, (box, label) in enumerate(zip(left_boxes, new_labels)):
        box.top    = r_top + i * (bh + GAP)
        box.height = bh
        set_bold_text(box, label, size=Pt(22), color=NAVY,
                      align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                      margin_left=Inches(0.15))
        print(f"    {box.name}: '{label}'  y={box.top/914400:.2f}\"  h={box.height/914400:.2f}\"")

    # Update title
    title = next(s for s in slide.shapes if s.name == "TextBox 1")
    set_bold_text(title, "Why design with Claude Code?",
                  size=Pt(30), color=NAVY,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.TOP,
                  margin_left=Inches(0))
    print("    Title updated")


# ── main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(PPTX)
    print(f"Loaded: {PPTX}  ({len(prs.slides)} slides)\n")

    print("Slide 2 — enlarging images...")
    modify_slide2_images(prs.slides[1])

    print("\nSlide 4 — replacing with two images...")
    modify_slide4_images(prs.slides[3])

    print("\nSlide 6 — fixing arrow overlap...")
    modify_slide6_arrow(prs.slides[5])

    print("\nSlide 7 — fixing layout + text...")
    modify_slide7_layout(prs.slides[6])

    prs.save(PPTX)
    print(f"\nSaved: {PPTX}")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""Apply requested edits to lightning-talk.pptx.

Run with: uv run --with python-pptx modify_lightning_talk.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "lightning-talk.pptx")

# ---- Color palette (same as build script) ----
NAVY       = RGBColor(0x1F, 0x36, 0x63)
BLUE       = RGBColor(0x00, 0x72, 0xBD)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE     = RGBColor(0xD4, 0x6A, 0x1A)
GREEN      = RGBColor(0x1A, 0x7A, 0x3C)
OLD_WAY_BG     = RGBColor(0xF8, 0xEE, 0xE4)
NEW_WAY_BG     = RGBColor(0xE4, 0xF4, 0xEA)
OLD_WAY_BORDER = RGBColor(0xD4, 0x8A, 0x50)
NEW_WAY_BORDER = RGBColor(0x4A, 0x9A, 0x5C)
BOX_BG     = RGBColor(0xEE, 0xF3, 0xF9)
CODE_BG    = RGBColor(0xF5, 0xF5, 0xF5)
CODE_BORDER= RGBColor(0xCC, 0xCC, 0xCC)
COMMENT_COLOR = RGBColor(0x5C, 0x7A, 0x3E)  # green for comments


# =========================================================
# Helpers
# =========================================================

def remove_shape(slide, shape):
    """Remove a shape from a slide via XML manipulation."""
    sp = shape._element
    sp.getparent().remove(sp)


def delete_slide(prs, index):
    """Delete a slide by index from the presentation."""
    slide_id_list = prs.slides._sldIdLst
    item = slide_id_list[index]
    r_id = item.attrib.get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    slide_id_list.remove(item)
    prs.part.drop_rel(r_id)


def add_rect(slide, x, y, w, h, fill=LIGHT_GRAY, border=MED_GRAY, border_pt=1.0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(border_pt)
    return shape


# =========================================================
# Slide modifiers
# =========================================================

def modify_slide2_projects(slide):
    """Remove bullet lists and the blue bottom note."""
    # TextBox 4 = left bullet list, TextBox 6 = right bullet list, TextBox 7 = blue bottom note
    to_remove = [s for s in slide.shapes if s.name in ("TextBox 4", "TextBox 6", "TextBox 7")]
    for s in to_remove:
        remove_shape(slide, s)
    print(f"  Slide 2: removed {len(to_remove)} shapes")


def modify_slide3_benchmarking(slide):
    """Replace old/new-way content with TOML design example."""
    # Remove everything except title (TextBox 1) and blue rule (Rectangle 2)
    to_remove = [s for s in slide.shapes if s.name not in ("TextBox 1", "Rectangle 2")]
    for s in to_remove:
        remove_shape(slide, s)

    # ---- OLD WAY box ----
    BX, BY, BW, BH = Inches(0.5), Inches(1.55), Inches(5.2), Inches(4.5)
    b = add_rect(slide, BX, BY, BW, BH, fill=OLD_WAY_BG, border=OLD_WAY_BORDER, border_pt=1.5)
    tf = b.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "THE OLD WAY"
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ORANGE

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = (
        '\n\u201cHow much do the most popular\n'
        'Python tools let users customize\n'
        'the format when writing TOML?\u201d\n\n'
        '\u2192 ask a developer\n'
        '\u2192 developer does research\n'
        '\u2192 comes back in a week\n'
        '\u2192 (or a month)'
    )
    r2.font.size = Pt(17); r2.font.color.rgb = DARK_GRAY

    # ---- Arrow ----
    arrow_tb = slide.shapes.add_textbox(Inches(5.85), Inches(3.3), Inches(0.6), Inches(0.6))
    arrow_tf = arrow_tb.text_frame
    ap = arrow_tf.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    ar = ap.add_run()
    ar.text = "\u27a1"
    ar.font.size = Pt(36); ar.font.color.rgb = NAVY

    # ---- NEW WAY box ----
    AX, AY, AW, AH = Inches(6.5), Inches(1.55), Inches(5.2), Inches(4.5)
    a = add_rect(slide, AX, AY, AW, AH, fill=NEW_WAY_BG, border=NEW_WAY_BORDER, border_pt=1.5)
    tf2 = a.text_frame
    tf2.word_wrap = True

    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "THE NEW WAY"
    r3.font.size = Pt(14); r3.font.bold = True; r3.font.color.rgb = GREEN

    p4 = tf2.add_paragraph()
    p4.alignment = PP_ALIGN.LEFT
    r4 = p4.add_run()
    r4.text = (
        '\nSame question \u2014 plus:\n'
        '\u201cHere are our design standards.\u201d\n\n'
        'Claude comes back with:\n'
        '\u2192 benchmarks across popular tools\n'
        '\u2192 recommended options to add\n'
        '\u2192 suggested names\n'
        '\u2192 immediately'
    )
    r4.font.size = Pt(17); r4.font.color.rgb = DARK_GRAY

    print("  Slide 3: replaced with TOML example")


def modify_slide4_usage(slide):
    """Remove text boxes (except title/rule), enlarge rescale image, add code snippet."""
    # Remove text boxes that aren't the title or rule
    to_remove = [s for s in slide.shapes if s.name in ("TextBox 4", "TextBox 5", "TextBox 6")]
    for s in to_remove:
        remove_shape(slide, s)

    # Remove Picture 8 (second smaller image)
    pic8 = [s for s in slide.shapes if s.name == "Picture 8"]
    for s in pic8:
        remove_shape(slide, s)

    # Enlarge Picture 7 (main rescale screenshot)
    for shape in slide.shapes:
        if shape.name == "Picture 7":
            shape.left  = Inches(0.35)
            shape.top   = Inches(1.25)
            shape.width = Inches(7.6)
            # Maintain aspect ratio: original was square (3931920 x 3931920 EMU)
            shape.height = Inches(7.6)
            # Clip to slide height: max useful height is about 6 inches
            shape.height = Inches(5.9)
            break

    # ---- Code snippet box on the right ----
    # From rescale_v5_sections_block.m, arguments block with block comments
    code_lines = [
        "arguments",
        "  %{",
        "  Input data, specified as a",
        "  vector, matrix, or N-D array.",
        "  If all elements of x are equal,",
        "  output contains NaN values.",
        "  %}",
        "  x        double",
        "",
        "  %{",
        "  Lower bound of the",
        "  target range.",
        "  %}",
        "  a (1,1)  double = 0",
        "",
        "  %{",
        "  Upper bound. If b < a,",
        "  output is reversed.",
        "  %}",
        "  b (1,1)  double = 1",
        "end",
    ]

    # Background box
    CX, CY, CW, CH = Inches(8.1), Inches(1.25), Inches(4.9), Inches(5.9)
    code_box = add_rect(slide, CX, CY, CW, CH, fill=CODE_BG, border=CODE_BORDER, border_pt=1.0)
    tf = code_box.text_frame
    tf.word_wrap = False

    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Courier New"
        run.font.size = Pt(12)
        # Color comments green
        if line.strip().startswith('%'):
            run.font.color.rgb = COMMENT_COLOR
        elif line == "arguments" or line == "end":
            run.font.color.rgb = BLUE
            run.font.bold = True
        else:
            run.font.color.rgb = DARK_GRAY

    print("  Slide 4: enlarged image, added code snippet")


def modify_slide6_configdata(slide):
    """Change 'ConfigData' to 'ConfigurationData' throughout the slide."""
    count = 0
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'ConfigData' in run.text:
                    run.text = run.text.replace('ConfigData', 'ConfigurationData')
                    count += 1
    print(f"  Slide 6: replaced ConfigData in {count} run(s)")


def modify_slide8_implications(slide):
    """Remove 'Memory and continuity' box and the bottom note."""
    to_remove = []
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        text = shape.text_frame.text
        if 'Memory and continuity' in text:
            to_remove.append(shape)
        elif 'nice-to-haves' in text:
            to_remove.append(shape)
    for s in to_remove:
        remove_shape(slide, s)
    print(f"  Slide 8: removed {len(to_remove)} shape(s)")


# =========================================================
# Main
# =========================================================

def main():
    prs = Presentation(PPTX_PATH)
    print(f"Loaded: {PPTX_PATH}  ({len(prs.slides)} slides)")

    # Apply changes using ORIGINAL slide indices (before any deletion)
    print("Modifying slide 2 (index 1)...")
    modify_slide2_projects(prs.slides[1])

    print("Modifying slide 3 (index 2)...")
    modify_slide3_benchmarking(prs.slides[2])

    print("Modifying slide 4 (index 3)...")
    modify_slide4_usage(prs.slides[3])

    print("Modifying slide 6 (index 5)...")
    modify_slide6_configdata(prs.slides[5])

    print("Modifying slide 8 (index 7)...")
    modify_slide8_implications(prs.slides[7])

    # Delete slide 7 (index 6) LAST to keep indices stable above
    print("Deleting slide 7 (index 6 = continuity slide)...")
    delete_slide(prs, 6)
    print(f"  Remaining slides: {len(prs.slides)}")

    prs.save(PPTX_PATH)
    print(f"\nSaved: {PPTX_PATH}")


if __name__ == "__main__":
    main()

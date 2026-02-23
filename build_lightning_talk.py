#!/usr/bin/env python3
"""Build lightning talk PPTX: 'Designing with a Partner: What AI Actually Changed'

Run with: uv run --with python-pptx build_lightning_talk.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUTPUT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "lightning-talk.pptx")

# ---- Color palette ----
NAVY        = RGBColor(0x1F, 0x36, 0x63)   # dark navy titles
BLUE        = RGBColor(0x00, 0x72, 0xBD)   # MathWorks accent blue
LIGHT_BLUE  = RGBColor(0x4D, 0xA6, 0xFF)   # lighter blue
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)   # body text
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)   # secondary text
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)  # light bg
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE      = RGBColor(0xD4, 0x6A, 0x1A)   # highlight
GREEN       = RGBColor(0x1A, 0x7A, 0x3C)   # success color
PLACEHOLDER_BG     = RGBColor(0xE4, 0xEF, 0xFF)
PLACEHOLDER_BORDER = RGBColor(0x88, 0xAA, 0xCC)
OLD_WAY_BG  = RGBColor(0xF8, 0xEE, 0xE4)  # warm orange-ish for "old way"
NEW_WAY_BG  = RGBColor(0xE4, 0xF4, 0xEA)  # green-ish for "new way"
OLD_WAY_BORDER = RGBColor(0xD4, 0x8A, 0x50)
NEW_WAY_BORDER = RGBColor(0x4A, 0x9A, 0x5C)
BOX_BG      = RGBColor(0xEE, 0xF3, 0xF9)  # neutral info box

# ---- EMU / sizing helpers ----
W = Inches(13.333)   # slide width
H = Inches(7.5)      # slide height


# =========================================================
# Helpers
# =========================================================

def new_slide(prs):
    """Add a blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])  # layout 6 = blank


def add_title(slide, text, y=Inches(0.28), h=Inches(0.85), size=Pt(30), color=NAVY):
    """Standard slide title."""
    tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.333), h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = size
    run.font.color.rgb = color
    return tb


def add_ruled_title(slide, text, size=Pt(28)):
    """Title + blue underline rule."""
    add_title(slide, text, size=size)
    rule = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.333), Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLUE
    rule.line.color.rgb = BLUE


def add_text(slide, text, x, y, w, h, size=Pt(18), color=DARK_GRAY,
             bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    """Plain text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    p.alignment = align
    return tb


def add_bullets(slide, lines, x, y, w, h, size=Pt(18), color=DARK_GRAY,
                bullet="•", spacing_after=Pt(4)):
    """Multi-line bullet text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"{bullet} {line}" if bullet else line
        run.font.size = size
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill=LIGHT_GRAY, border=MED_GRAY, border_pt=1.0):
    """Filled rectangle shape."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(border_pt)
    return shape


def add_labeled_rect(slide, label, x, y, w, h, fill=LIGHT_GRAY, border=MED_GRAY,
                     label_size=Pt(16), label_color=DARK_GRAY, bold=False):
    """Rectangle with centered text label."""
    shape = add_rect(slide, x, y, w, h, fill=fill, border=border)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = label_size
    run.font.color.rgb = label_color
    run.font.bold = bold
    return shape


def add_placeholder(slide, label, x, y, w, h):
    """Blue-tinted placeholder box for screenshots."""
    shape = add_rect(slide, x, y, w, h, fill=PLACEHOLDER_BG, border=PLACEHOLDER_BORDER, border_pt=1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x44, 0x66, 0x99)
    run.font.italic = True
    return shape


def add_arrow_text(slide, text, x, y, size=Pt(28), color=NAVY):
    """Just a text arrow glyph."""
    tb = slide.shapes.add_textbox(x, y, Inches(0.6), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    return tb


def set_notes(slide, text):
    """Set speaker notes text."""
    slide.notes_slide.notes_text_frame.text = text


# =========================================================
# Slide builders
# =========================================================

def slide1_hook(prs):
    """Hook slide: the central claim."""
    s = new_slide(prs)

    # Main quote - two lines, big
    tb = slide.shapes.add_textbox if False else s.shapes.add_textbox(
        Inches(1.0), Inches(1.2), Inches(11.333), Inches(2.8))
    tf = tb.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "\u201cI didn\u2019t use AI to write code."
    r1.font.size = Pt(40)
    r1.font.color.rgb = NAVY
    r1.font.bold = True

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "I used it to do design.\u201d"
    r2.font.size = Pt(40)
    r2.font.color.rgb = BLUE
    r2.font.bold = True

    # Byline
    tb2 = s.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.333), Inches(0.6))
    tf2 = tb2.text_frame
    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Michelle Hirsch  \u00b7  MathWorks"
    r3.font.size = Pt(20)
    r3.font.color.rgb = MED_GRAY

    # Framing question
    tb3 = s.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10.333), Inches(1.5))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p4 = tf3.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    r4.text = "What does it mean for AI to be a design partner \u2014 not a code generator?"
    r4.font.size = Pt(21)
    r4.font.color.rgb = DARK_GRAY
    r4.font.italic = True

    set_notes(s, """\
• Introduce yourself: Michelle Hirsch, MathWorks — you work on MATLAB tooling
• Frame the surprise: most of what you'll hear about AI in engineering is code generation
• Your experience was different — it was design collaboration
• Two concrete projects as evidence — not theory, not demos, actual work products
• This is a practitioner report: honest, specific, no hype""")
    return s


def slide2_projects(prs):
    """Two concrete projects."""
    s = new_slide(prs)
    add_ruled_title(s, "Two Concrete Projects", size=Pt(30))

    # --- Left column: Project 1 ---
    MID = Inches(6.9)  # midpoint between columns
    LEFT_X = Inches(0.5)
    COL_W = Inches(5.9)

    # Header box - Project 1
    add_labeled_rect(s, "MATLAB Documentation Framework",
                     LEFT_X, Inches(1.35), COL_W, Inches(0.55),
                     fill=BLUE, border=BLUE, label_size=Pt(16), label_color=WHITE, bold=True)

    add_bullets(s, [
        ".m help comments \u2192 rich HTML  (like mkdocstrings for Python)",
        "720-line living spec with rationale",
        "40+ progressive sample files",
        "Working prototype",
        "Pure design work: grammar, decisions, tradeoffs",
    ], LEFT_X + Inches(0.1), Inches(2.0), COL_W - Inches(0.1), Inches(3.5),
               size=Pt(16), color=DARK_GRAY, bullet="\u2022")

    # --- Right column: Project 2 ---
    RIGHT_X = Inches(7.0)

    add_labeled_rect(s, "Config File Readers & Writers for MATLAB",
                     RIGHT_X, Inches(1.35), COL_W, Inches(0.55),
                     fill=NAVY, border=NAVY, label_size=Pt(16), label_color=WHITE, bold=True)

    add_bullets(s, [
        "YAML, TOML, INI, and JSON formats",
        "Custom typed data classes (not plain struct)",
        "Shared ConfigData base class",
        "Consistent API across all 4 formats",
        "4 formats designed together \u2014 drove cross-format consistency",
    ], RIGHT_X + Inches(0.1), Inches(2.0), COL_W - Inches(0.1), Inches(3.5),
               size=Pt(16), color=DARK_GRAY, bullet="\u2022")

    # Bottom label
    add_text(s, "Both were primarily design challenges \u2014 not coding challenges.",
             Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.55),
             size=Pt(17), color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    set_notes(s, """\
• Project 1: The MATLAB documentation framework
  - The idea: .m help comments should render as rich HTML when you call doc()
  - Like what mkdocstrings does for Python — but for MATLAB
  - The design challenge: what grammar? what sections? how do arguments get documented?
  - Work products: 720-line living spec, 40+ sample .m files (7 versions of rescale, 6 of Sensor class), working prototype
  - This was almost entirely design work — grammar design, tradeoff analysis, decision records

• Project 2: Config file I/O for MATLAB
  - YAML, TOML, INI, JSON readers and writers
  - The interesting design decisions: should config data be struct (easy but weak) or typed classes?
  - Went with typed classes: YAMLData, TOMLData, INIData, JSONData
  - Discovered they share structure → extracted shared ConfigData base class
  - Designed all 4 simultaneously → forced consistent API across formats

• Neither project was primarily about code — they were design collaborations""")
    return s


def slide3_benchmarking(prs):
    """Benchmarking at the speed of thought."""
    s = new_slide(prs)
    add_ruled_title(s, "Benchmarking No Longer Interrupts the Design Conversation")

    # Before box
    BX = Inches(0.5)
    BY = Inches(1.55)
    BW = Inches(5.2)
    BH = Inches(3.6)
    b = add_rect(s, BX, BY, BW, BH, fill=OLD_WAY_BG, border=OLD_WAY_BORDER, border_pt=1.5)
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "THE OLD WAY"; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = ORANGE

    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = ('\n\u201cI need to check how\n'
               'mkdocstrings handles\n'
               'argument documentation\u2026\u201d\n\n'
               '\u2192 leave the conversation\n'
               '\u2192 open browser, research\n'
               '\u2192 come back\n'
               '\u2192 lost the thread')
    r2.font.size = Pt(16); r2.font.color.rgb = DARK_GRAY

    # Arrow between
    add_arrow_text(s, "\u27a1", Inches(5.85), Inches(3.0), size=Pt(36), color=NAVY)

    # After box
    AX = Inches(6.5)
    AY = Inches(1.55)
    AW = Inches(5.2)
    AH = Inches(3.6)
    a = add_rect(s, AX, AY, AW, AH, fill=NEW_WAY_BG, border=NEW_WAY_BORDER, border_pt=1.5)
    tf2 = a.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = "THE NEW WAY"; r3.font.size = Pt(14); r3.font.bold = True
    r3.font.color.rgb = GREEN

    p4 = tf2.add_paragraph(); p4.alignment = PP_ALIGN.LEFT
    r4 = p4.add_run()
    r4.text = ('\n\u201cHow does mkdocstrings\n'
               'handle this?\u201d\n\n'
               '\u2192 get answer\n'
               '\u2192 keep designing\n\n'
               'No context switch.\n'
               'No lost thread.')
    r4.font.size = Pt(16); r4.font.color.rgb = DARK_GRAY

    # Concrete example footer
    ex_box = add_rect(s, Inches(0.5), Inches(5.35), Inches(12.333), Inches(0.85),
                      fill=BOX_BG, border=BLUE, border_pt=1.0)
    tf3 = ex_box.text_frame
    tf3.word_wrap = True
    p5 = tf3.paragraphs[0]
    r5 = p5.add_run()
    r5.text = ("Concrete: Deciding argument-doc syntax \u2192 immediately checked Sphinx, mkdocstrings, Javadoc. "
               "Made a decision and kept going. For config formats: checked TOML spec edge cases, "
               "Python library behavior, real user patterns \u2014 all inline.")
    r5.font.size = Pt(14); r5.font.color.rgb = DARK_GRAY; r5.font.italic = True

    set_notes(s, """\
• The old way: any benchmarking required leaving the conversation
  - Open browser, read docs, come back — broke the design flow
  - "I'll look that up later" often meant never
  - Context switch is expensive

• The new way: ask inline, get answer, continue
  - No interruption to the design conversation
  - Benchmarking became part of the conversation, not a detour

• Concrete example from the doc framework:
  - Decision: how should argument short-descriptions work in our grammar?
  - Immediately checked how Sphinx, mkdocstrings, and Javadoc handle it
  - Compared approaches, chose the best for our use case, moved on
  - The whole thing took maybe 3 minutes inside the existing conversation

• Config format examples:
  - TOML spec edge cases (e.g., how does multiline string escaping work?)
  - How Python's tomllib handles ambiguous cases
  - What options real users actually need vs. rarely use
  - All checked while designing the API — zero context switch

• The point: Benchmarking is no longer a separate research phase.
  It's integrated into the design conversation.""")
    return s


def slide4_usage(prs):
    """Designing through usage."""
    s = new_slide(prs)
    add_ruled_title(s, "Designs Get Better Through Usage, Not by Staring at Specs")

    # Left: placeholder for screenshot
    add_placeholder(s, "[SCREENSHOT: rescale v0 (bare arguments block only)\nvs. rescale v5 (rich documentation)\nas rendered HTML pages]",
                    Inches(0.5), Inches(1.4), Inches(6.8), Inches(4.5))

    # Right: key points
    add_text(s, "Validated by building, not speccing:", Inches(7.55), Inches(1.45),
             Inches(5.3), Inches(0.5), size=Pt(18), color=NAVY, bold=True)

    add_bullets(s, [
        "7 progressive versions of rescale()",
        "6 versions of a Sensor class",
        "Each actually rendered as HTML",
        "Problems obvious in rendered output, invisible in spec",
        "Low prototyping cost = willing to iterate",
    ], Inches(7.55), Inches(2.1), Inches(5.3), Inches(3.3),
               size=Pt(17), color=DARK_GRAY)

    add_text(s, "Also: working config-format implementations let us test\ndesign decisions against real data and real edge cases.",
             Inches(0.5), Inches(6.15), Inches(12.333), Inches(0.85),
             size=Pt(15), color=MED_GRAY, italic=True)

    set_notes(s, """\
• The grammar was validated by building real examples, not by reviewing the spec

• 7 progressively documented versions of rescale():
  - v0: bare function, just the arguments block (auto-generated minimal page)
  - v1: add first-line synopsis only
  - v2: add trailing comments on argument lines (short descriptions)
  - v3: add block comments before argument lines (longer descriptions)
  - v4: add a full ## Input Arguments section
  - v5: full documentation with Examples and Tips
  - v6/v7: testing edge cases and class documentation

• Also 6 versions of a Sensor class — validating class-specific features
  (properties, methods, events)

• All of these were actually rendered as HTML using the working prototype

• Key insight: things that look fine in a spec become obviously wrong
  when you see real rendered output
  - Visual hierarchy issues
  - Awkward section headings
  - Cases where "section wins" vs "merge" matters

• For the config formats: working implementations meant design decisions
  were validated against real YAML/TOML/INI/JSON files, not hypothetical ones

• The point: rapid prototyping at low cost changes what you're willing to try.
  You don't just design differently — you validate differently.""")
    return s


def slide5_rationale(prs):
    """Capturing rationale."""
    s = new_slide(prs)
    add_ruled_title(s, "The Spec Captures Not Just Decisions\u2014But Why")

    # Left: screenshot placeholder
    add_placeholder(s, "[SCREENSHOT: spec section showing\n\u2018Rationale\u2019 + \u2018Alternatives considered\u2019\nblocks, with Syntax Description Gap\ncallout highlighted]",
                    Inches(0.5), Inches(1.4), Inches(6.3), Inches(4.0))

    # Right: explanation
    add_text(s, "Example: The Syntax Description Gap",
             Inches(7.0), Inches(1.45), Inches(5.8), Inches(0.5),
             size=Pt(18), color=NAVY, bold=True)

    add_text(s, ("A structural tension: if you add descriptions to calling forms, "
                 "you implicitly take full ownership of the syntax list.\n\n"
                 "Documented: why \u201csection wins entirely\u201d instead of a merge model."),
             Inches(7.0), Inches(2.1), Inches(5.8), Inches(2.0),
             size=Pt(16), color=DARK_GRAY, wrap=True)

    # Bottom box: the point
    bottom = add_rect(s, Inches(0.5), Inches(5.6), Inches(12.333), Inches(1.0),
                      fill=BOX_BG, border=BLUE)
    tf = bottom.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("Design reasoning is normally lost. Engineers inherit decisions with no context. "
              "Here, the rationale lives in the document \u2014 alongside the decision.")
    r.font.size = Pt(16); r.font.color.rgb = DARK_GRAY

    set_notes(s, """\
• The spec has explicit "Rationale" and "Alternatives considered" sections throughout — not just decisions

• Specific example: The "Syntax Description Gap"
  - Structural tension in the grammar design
  - If you want to add descriptions to individual calling forms (syntax variants),
    you implicitly take ownership of listing ALL calling forms
  - The auto-generated list from the arguments block can't be merged with a user-defined list
  - Decision: "section wins entirely" — if you define a ## Syntax section, it completely
    replaces the auto-generated one
  - Rationale: documented in the spec with the tradeoffs explicitly stated

• Why does this matter?
  - Normally, design reasoning is lost the moment the meeting ends
  - Engineers who come later have no idea WHY things are the way they are
  - They either reverse decisions accidentally, or fear changing anything
  - Having the rationale documented changes what future maintainers can do

• Having AI in the loop made capturing rationale feel natural, not overhead
  - Articulating "here's why we chose this" became part of the conversation
  - The AI would sometimes surface considerations I hadn't thought of, which made
    the rationale richer

• The point: AI doesn't just help you make decisions — it helps you articulate
  and document the reasoning behind them.""")
    return s


def slide6_brave(prs):
    """Brave design through reduced activation energy."""
    s = new_slide(prs)
    add_ruled_title(s, "Lower Cost Changed What I Was Willing to Attempt")

    # Evolution diagram: struct → 4 types → base class
    # Row 1: "Without AI" scenario
    add_text(s, "Without AI:", Inches(0.5), Inches(1.45), Inches(2.0), Inches(0.4),
             size=Pt(16), color=MED_GRAY, bold=True)
    add_labeled_rect(s, "struct\n(plain dict)", Inches(2.6), Inches(1.35),
                     Inches(2.2), Inches(0.75), fill=OLD_WAY_BG, border=OLD_WAY_BORDER,
                     label_size=Pt(16), label_color=DARK_GRAY)
    add_text(s, "\u2190  too daunting to design proper types",
             Inches(4.9), Inches(1.55), Inches(7.0), Inches(0.4),
             size=Pt(15), color=ORANGE, italic=True)

    # Row 2: "With AI" — 4 custom types
    add_text(s, "With AI:", Inches(0.5), Inches(2.45), Inches(2.0), Inches(0.4),
             size=Pt(16), color=NAVY, bold=True)
    for i, (name, clr) in enumerate([
            ("YAMLData", BLUE),
            ("TOMLData", NAVY),
            ("INIData",  RGBColor(0x2E, 0x7D, 0x32)),
            ("JSONData", RGBColor(0x6A, 0x1A, 0x9A)),
    ]):
        x = Inches(2.5 + i * 2.5)
        add_labeled_rect(s, name, x, Inches(2.35), Inches(2.1), Inches(0.65),
                         fill=clr, border=clr, label_size=Pt(16),
                         label_color=WHITE, bold=True)

    # Arrow down
    add_arrow_text(s, "\u2193", Inches(6.3), Inches(3.15), size=Pt(28), color=NAVY)
    add_text(s, "Shared structure discovered \u2192 extracted base class",
             Inches(2.5), Inches(3.15), Inches(8.5), Inches(0.4),
             size=Pt(15), color=DARK_GRAY, italic=True)

    # Row 3: base class
    add_labeled_rect(s, "ConfigData  (shared base class)",
                     Inches(3.5), Inches(3.7), Inches(5.8), Inches(0.7),
                     fill=BOX_BG, border=BLUE, label_size=Pt(18), label_color=NAVY, bold=True)

    # 4-format insight
    cross_box = add_rect(s, Inches(0.5), Inches(4.7), Inches(12.333), Inches(0.9),
                         fill=NEW_WAY_BG, border=NEW_WAY_BORDER, border_pt=1.2)
    tf = cross_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("Designing all 4 formats simultaneously drove cross-format consistency "
              "that would never have happened if each was designed in isolation.")
    r.font.size = Pt(16); r.font.color.rgb = DARK_GRAY

    # Bottom punch line
    add_text(s, "AI doesn\u2019t just lower effort per task \u2014 it changes what you\u2019re willing to attempt.",
             Inches(0.5), Inches(5.85), Inches(12.333), Inches(0.55),
             size=Pt(18), color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    set_notes(s, """\
• Without AI: I would have represented config data as plain struct
  - MATLAB struct is familiar, easy, zero design effort
  - Designing custom types felt too daunting: too many decisions, too much work
  - The result would have been: functional but weak

• With AI: We designed proper typed classes
  - YAMLData, TOMLData, INIData, JSONData — each a proper typed class
  - Conversations: what fields? what methods? what do you do with nested structures?
  - Prototyped quickly enough that we could try things and see if they worked

• Discovered shared structure → extracted ConfigData base class
  - All 4 types had common fields and methods
  - Natural refactoring: extract to a shared base class
  - This wouldn't have happened if types started as structs

• The cross-format consistency point:
  - Working on all 4 simultaneously meant decisions had to be consistent
  - "How should YAML handle nested maps?" → immediately applied the same pattern to TOML, INI, JSON
  - The result: a consistent, coherent API across 4 formats
  - Isolation would have produced 4 slightly inconsistent designs

• This is the "activation energy" insight:
  - Not just "it's faster to do the same thing"
  - It changes WHAT you attempt
  - Better designs emerge when the cost of ambition is lower""")
    return s


def slide7_continuity(prs):
    """Continuity and memory."""
    s = new_slide(prs)
    add_ruled_title(s, "Design Is a Conversation That Spans Many Sessions")

    # Timeline-ish visual: three session boxes connected by a line
    LINE_Y = Inches(3.3)

    # Connecting line
    line_shape = add_rect(s, Inches(1.2), LINE_Y, Inches(10.9), Pt(3),
                          fill=BLUE, border=BLUE, border_pt=0)

    session_data = [
        ("Session A\n(weeks ago)", "Grammar design,\nspec started", Inches(1.2)),
        ("Session B\n(last week)", "Revise after\nprototype testing", Inches(5.2)),
        ("Session C\n(today)", "Fill gaps,\nclean up spec", Inches(9.2)),
    ]

    for label, sublabel, x in session_data:
        # Circle-ish blob — use a rectangle, centered on the line
        blob = add_rect(s, x - Inches(0.07), LINE_Y - Inches(0.12),
                        Inches(0.28), Inches(0.28),
                        fill=BLUE, border=BLUE)
        # Session box above the line
        box = add_rect(s, x - Inches(0.7), Inches(1.45), Inches(1.9), Inches(1.5),
                       fill=BOX_BG, border=BLUE)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = NAVY
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sublabel; r2.font.size = Pt(13)
        r2.font.color.rgb = DARK_GRAY

        # Connector line from box to timeline
        connector = add_rect(s, x - Inches(0.015), Inches(2.95), Pt(3), Inches(0.35),
                             fill=BLUE, border=BLUE)

    # Key points below timeline
    add_bullets(s, [
        "Come back after days or weeks. Share the spec. Pick up exactly where you left off.",
        "AI reconstructs context from the spec and sample files.",
        "Also enables cleanup: revisit spec after implementation to reconcile and fill gaps.",
    ], Inches(0.7), Inches(3.85), Inches(11.9), Inches(2.0),
               size=Pt(17), color=DARK_GRAY)

    add_text(s, "Design becomes an ongoing iterative process, not a one-time event.",
             Inches(0.5), Inches(6.15), Inches(12.333), Inches(0.55),
             size=Pt(18), color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    set_notes(s, """\
• These projects lasted weeks to months — not single sessions
• The spec as "shared memory" artifact:
  - Start each new session by pointing the AI to the spec and sample files
  - It reconstructs the context and continues cleanly
  - No need to re-explain decisions made weeks ago

• This makes design genuinely iterative:
  - Not just "iterate within a session" but iterate across time
  - Come back after seeing real usage, come back after implementation reveals new questions

• The cleanup use case:
  - After building the prototype, I went back to the spec
  - Found gaps where implementation diverged from spec
  - Filled in sections that were vague, reconciled where spec was wrong
  - The spec stayed current with reality

• This is a quality story too:
  - Without continuity, specs get stale immediately after they're written
  - With continuity, the spec remains a living document

• Quick slide — don't dwell — but it's worth naming as a real benefit""")
    return s


def slide8_implications(prs):
    """Implications for the audience."""
    s = new_slide(prs)
    add_ruled_title(s, "You Need These Things Too", size=Pt(32))

    # Left: capability list
    caps = [
        ("Instant benchmarking",         "Access to specs, literature, tool comparisons in-context"),
        ("Rapid prototyping",             "Try ideas with their own tools and data at low cost"),
        ("Help articulating reasoning",   "Not just answers — help documenting the why"),
        ("Support for ambitious designs", "Attempt goals they\u2019d otherwise skip"),
        ("Memory and continuity",         "Pick up long projects without re-explaining context"),
    ]

    for i, (cap, desc) in enumerate(caps):
        y = Inches(1.45 + i * 0.96)
        box = add_rect(s, Inches(0.5), y, Inches(7.7), Inches(0.82),
                       fill=BOX_BG, border=BLUE, border_pt=0.8)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = cap + ":  "; r1.font.size = Pt(15); r1.font.bold = True
        r1.font.color.rgb = NAVY
        r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(15)
        r2.font.color.rgb = DARK_GRAY

    # Right: the question
    q_box = add_rect(s, Inches(8.45), Inches(1.45), Inches(4.4), Inches(4.9),
                     fill=RGBColor(0xFE, 0xF9, 0xE7), border=ORANGE, border_pt=1.5)
    tf = q_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "\u201cWhat would change about your design process if your AI could do these things for you?\u201d"
    r.font.size = Pt(17); r.font.color.rgb = DARK_GRAY; r.font.italic = True

    add_text(s, "These aren\u2019t nice-to-haves. They change what scientists design and how ambitiously they design it.",
             Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.55),
             size=Pt(16), color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    set_notes(s, """\
• This is the meta-lesson: you are building AI tools for scientists who do design work
• The things that changed my design process are the things scientists need:

1. Instant benchmarking
   - Literature access, precedent checking, tool comparisons — in-context, while designing
   - Not a separate research phase

2. Rapid prototyping
   - Try ideas with their actual data and tools
   - Low cost = willing to try things that might not work

3. Help articulating reasoning
   - Not just "give me the answer"
   - Help saying WHY — designing better, documenting better

4. Support for ambitious designs
   - Lower activation energy for harder problems
   - Better designs emerge when the cost of ambition is lower

5. Memory and continuity
   - Long projects don't restart from scratch each session
   - The context lives in the artifacts, not in human memory

• Turn the question back to the audience:
  - You've heard what changed for me
  - What would change for your users?
  - What would change for YOU, as designers of these tools?

• The prompt at the end is for reflection — let it land""")
    return s


def slide9_close(prs):
    """Closing slide."""
    s = new_slide(prs)

    # Main claim — centered, large
    tb = s.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\u201cThe quality bar is different"
    r.font.size = Pt(42); r.font.bold = True; r.font.color.rgb = NAVY

    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "when you have a design partner.\u201d"
    r2.font.size = Pt(42); r2.font.bold = True; r2.font.color.rgb = BLUE

    # Blue divider
    div = add_rect(s, Inches(3.0), Inches(3.85), Inches(7.333), Pt(3),
                   fill=BLUE, border=BLUE)

    # Invitation
    tb2 = s.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(11.333), Inches(1.0))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p3 = tf2.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Try it on your next design decision.\u2003Not your next coding task."
    r3.font.size = Pt(22); r3.font.color.rgb = DARK_GRAY

    # Name at bottom
    tb3 = s.shapes.add_textbox(Inches(1.0), Inches(6.3), Inches(11.333), Inches(0.6))
    tf3 = tb3.text_frame
    p4 = tf3.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    r4.text = "Michelle Hirsch  \u00b7  MathWorks"
    r4.font.size = Pt(18); r4.font.color.rgb = MED_GRAY

    set_notes(s, """\
• Close on the single sentence: "The quality bar is different when you have a design partner."
• Not because AI does the work for you
• But because you attempt things you wouldn't attempt alone
• You validate assumptions you'd otherwise skip
• You document reasoning you'd otherwise lose
• You work on 4 formats instead of 1

• The invitation is genuine: try it on a design problem, not a coding task
  - Pick a real design decision you're facing
  - Use it to benchmark options, prototype quickly, articulate the rationale
  - See what changes

• If you've been using AI for code generation, this is a different mode
  - It requires treating the AI as a collaborator, not a tool
  - The spec/artifact becomes the shared memory
  - The conversation drives the design

• Thank you""")
    return s


# =========================================================
# Main
# =========================================================

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("Building slides...")
    slide1_hook(prs)
    print("  Slide 1: Hook")
    slide2_projects(prs)
    print("  Slide 2: Projects")
    slide3_benchmarking(prs)
    print("  Slide 3: Benchmarking")
    slide4_usage(prs)
    print("  Slide 4: Designing Through Usage")
    slide5_rationale(prs)
    print("  Slide 5: Capturing Rationale")
    slide6_brave(prs)
    print("  Slide 6: Brave Design")
    slide7_continuity(prs)
    print("  Slide 7: Continuity")
    slide8_implications(prs)
    print("  Slide 8: Implications")
    slide9_close(prs)
    print("  Slide 9: Close")

    prs.save(OUTPUT_PATH)
    print(f"\nSaved: {OUTPUT_PATH}")


if __name__ == "__main__":
    build_presentation()
